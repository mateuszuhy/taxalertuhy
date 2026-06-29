import os
from pptx import Presentation


def create_ppt(data):

    # 🔥 HARD FIX: ensure output directory exists
    os.makedirs("output", exist_ok=True)

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "UHY Tax Alert V3.3"

    # Lead news
    for n in data.get("lead", []):

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = n.get("title", "No title")
        slide.placeholders[1].text = "\n".join(n.get("summary", []))

    # Standard news
    for n in data.get("standard", []):

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = n.get("title", "No title")
        slide.placeholders[1].text = "\n".join(n.get("summary", []))

    file_path = "output/tax_alert.pptx"
    prs.save(file_path)

    return file_path
